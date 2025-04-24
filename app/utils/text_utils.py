#!/usr/bin/env python3
"""
Text utility functions for the Mail Analysis API.
"""

import re
from typing import Optional, Union
from loguru import logger
import json
import io
from typing import Dict, List, Any
import base64

# Try importing pandas, which is the best library for Excel handling
try:
    import pandas as pd
    pandas_available = True
except ImportError:
    pandas_available = False
    logger.warning("pandas library not available, Excel conversion will be limited")

try:
    import html2text
    html2text_available = True
except ImportError:
    html2text_available = False
    logger.warning("html2text library not available, using fallback HTML to Markdown conversion")

# Try importing PDF libraries
try:
    import pymupdf4llm
    import fitz  # PyMuPDF

    pymupdf4llm_available = True
except ImportError:
    logger.warning("PDF libraries not available, PDF conversion will be limited")

# Try importing PowerPoint libraries
try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False
    logger.warning("python-pptx library not available, PowerPoint conversion will be limited")

# Try importing Word document libraries
try:
    import docx
    docx_available = True
except ImportError:
    docx_available = False
    try:
        # Try an alternative library for older .doc files
        import textract
        textract_available = True
    except ImportError:
        textract_available = False
        logger.warning("Word document processing libraries not available, Word conversion will be limited")

def html_to_markdown(html_content: str) -> str:
    """
    Convert HTML content to Markdown format.
    
    Args:
        html_content: HTML content to convert
        
    Returns:
        Markdown formatted text
    """
    if not html_content:
        return ""
    
    try:
        if html2text_available:
            # Use html2text library for conversion
            converter = html2text.HTML2Text()
            converter.ignore_links = False
            converter.ignore_images = False
            converter.ignore_tables = False
            converter.body_width = 0  # Don't wrap text
            return converter.handle(html_content)
        else:
            # Fallback simple HTML to Markdown conversion
            markdown = html_content
            
            # Remove HTML doctype and meta tags
            markdown = re.sub(r'<!DOCTYPE.*?>', '', markdown, flags=re.DOTALL)
            markdown = re.sub(r'<head>.*?</head>', '', markdown, flags=re.DOTALL)
            
            # Convert common HTML elements to Markdown
            # Headers
            markdown = re.sub(r'<h1>(.*?)</h1>', r'# \1', markdown)
            markdown = re.sub(r'<h2>(.*?)</h2>', r'## \1', markdown)
            markdown = re.sub(r'<h3>(.*?)</h3>', r'### \1', markdown)
            
            # Lists
            markdown = re.sub(r'<ul>(.*?)</ul>', r'\1', markdown, flags=re.DOTALL)
            markdown = re.sub(r'<li>(.*?)</li>', r'- \1\n', markdown)
            
            # Links
            markdown = re.sub(r'<a href="(.*?)">(.*?)</a>', r'[\2](\1)', markdown)
            
            # Bold and Italic
            markdown = re.sub(r'<strong>(.*?)</strong>', r'**\1**', markdown)
            markdown = re.sub(r'<b>(.*?)</b>', r'**\1**', markdown)
            markdown = re.sub(r'<em>(.*?)</em>', r'*\1*', markdown)
            markdown = re.sub(r'<i>(.*?)</i>', r'*\1*', markdown)
            
            # Paragraphs and line breaks
            markdown = re.sub(r'<p>(.*?)</p>', r'\1\n\n', markdown, flags=re.DOTALL)
            markdown = re.sub(r'<br\s*/?>', r'\n', markdown)
            
            # Tables - simplified conversion (without alignment)
            markdown = re.sub(r'<table>(.*?)</table>', r'\1', markdown, flags=re.DOTALL)
            markdown = re.sub(r'<tr>(.*?)</tr>', r'\1\n', markdown, flags=re.DOTALL)
            markdown = re.sub(r'<th>(.*?)</th>', r'| \1 ', markdown)
            markdown = re.sub(r'<td>(.*?)</td>', r'| \1 ', markdown)
            
            # Remove any remaining HTML tags
            markdown = re.sub(r'<.*?>', '', markdown)
            
            # Normalize newlines and whitespace
            markdown = re.sub(r'\n{3,}', r'\n\n', markdown)
            markdown = markdown.strip()
            
            return markdown
            
    except Exception as e:
        logger.error(f"Error converting HTML to Markdown: {str(e)}")
        # Return raw HTML if conversion fails
        return html_content

def clean_text(text: str, remove_urls: bool = False, remove_emails: bool = False) -> str:
    """
    Clean text by removing unwanted elements.
    
    Args:
        text: Text to clean
        remove_urls: Flag to remove URLs
        remove_emails: Flag to remove email addresses
        
    Returns:
        Cleaned text
    """
    if not text:
        return ""
    
    # Remove URLs if requested
    if remove_urls:
        text = re.sub(r'https?://\S+', '', text)
        
    # Remove email addresses if requested
    if remove_emails:
        text = re.sub(r'\S+@\S+\.\S+', '', text)
    
    # Remove multiple newlines
    text = re.sub(r'\n{3,}', r'\n\n', text)
    
    # Remove extra whitespace
    text = re.sub(r'\s{2,}', ' ', text)
    
    return text.strip()

def convert_excel_to_json(excel_file: str) -> str:
    """
    Convert Excel content to JSON format.
    
    Args:
        excel_file: Path to the Excel file
        
    Returns:
        JSON string representation of the Excel data
    """
    if not excel_file:
        return ""
    
    try:
        # If content is a JSON string, just return it
        if isinstance(excel_file, str) and (excel_file.startswith("{") or excel_file.startswith("[")):
            return excel_file
            
        if pandas_available:
            # Read all sheets from the Excel file
            result = {}
            excel_data = pd.read_excel(excel_file, sheet_name=None)
            
            # Convert each sheet to JSON
            for sheet_name, df in excel_data.items():
                # Convert datetime columns to strings
                for col in df.select_dtypes(include=['datetime64']).columns:
                    df[col] = df[col].dt.strftime('%Y-%m-%d %H:%M:%S')
                
                # Replace NaN values with None for proper JSON serialization
                df = df.where(pd.notna(df), None)
                
                # Convert to list of dictionaries (records)
                sheet_data = df.to_dict(orient='records')
                
                # Process any remaining Timestamp objects
                processed_data = []
                for record in sheet_data:
                    processed_record = {}
                    for key, value in record.items():
                        # Convert any pandas Timestamp objects to strings
                        if hasattr(value, 'strftime'):
                            processed_record[str(key)] = value.strftime('%Y-%m-%d %H:%M:%S')
                        else:
                            processed_record[str(key)] = value
                    processed_data.append(processed_record)
                    
                result[str(sheet_name)] = processed_data
            
            # Custom JSON encoder for any other date/time types
            class CustomJSONEncoder(json.JSONEncoder):
                def default(self, obj):
                    if hasattr(obj, 'strftime'):
                        return obj.strftime('%Y-%m-%d %H:%M:%S')
                    return super().default(obj)
            
            # Convert the final result to JSON using the custom encoder
            return json.dumps(result, ensure_ascii=False, cls=CustomJSONEncoder)
        else:
            # Fallback for when pandas isn't available
            logger.warning("pandas not available for Excel conversion, returning simple text")
            return "Excel content couldn't be converted without pandas library."
    except Exception as e:
        logger.error(f"Error converting Excel to JSON: {str(e)}")
        # Return a simplified error message in JSON format
        return json.dumps({
            "error": f"Failed to convert Excel file: {str(e)}"
        })

def convert_pdf_to_markdown(pdf_file:str) -> str:
    """
    Convert PDF content to Markdown text format.
    
    Args:
        pdf_content: PDF file content as bytes or base64 encoded string
        
    Returns:
        Markdown formatted text extracted from the PDF
    """
 
    
    # Step 2: Open PDF bytes as a PyMuPDF Document
    pdf_document = fitz.open(filename=pdf_file, filetype="pdf")

    markdown_content = pymupdf4llm.to_markdown(pdf_document)

    return markdown_content


def convert_ppt_to_markdown(ppt_file: str) -> str:
    """
    Convert PowerPoint content to Markdown text format.
    
    Args:
        ppt_file: Full path to the PowerPoint file
        
    Returns:
        Markdown formatted text extracted from the PowerPoint
    """
    if not ppt_file:
        return ""
    
    try:
        extracted_text = ""
        
        if pptx_available:
            # Use python-pptx for extraction - open directly from file path
            presentation = Presentation(ppt_file)
            
            # Extract text from each slide
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract title if present
                if slide.shapes.title:
                    slide_text.append(f"### {slide.shapes.title.text}")
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if this is the title we already added
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text)
                
                # Combine all text from this slide
                if slide_text:
                    extracted_text += f"## Slide {slide_num + 1}\n\n"
                    extracted_text += "\n\n".join(slide_text) + "\n\n"
                else:
                    extracted_text += f"## Slide {slide_num + 1}\n\n*No text content*\n\n"
        else:
            return "PowerPoint text extraction failed. The python-pptx library is not installed."
        
        # Clean up the text
        extracted_text = clean_text(extracted_text)
        
        return extracted_text
            
    except Exception as e:
        logger.error(f"Error converting PowerPoint to Markdown: {str(e)}")
        return f"Error extracting text from PowerPoint: {str(e)}"

def convert_word_to_markdown(word_file: str) -> str:
    """
    Convert Word document content to Markdown text format.
    
    Args:
        word_file: Path to the Word document file
        
    Returns:
        Markdown formatted text extracted from the Word document
    """
    if not word_file:
        return ""
    
    try:
        extracted_text = ""
        
        # Check if it's a .docx file by examining the file extension
        is_docx = word_file.lower().endswith('.docx')
        
        if docx_available and is_docx:
            # Use python-docx for .docx files
            doc = docx.Document(word_file)
            
            # Process document structure and convert to markdown
            
            # Extract title if available
            if doc.paragraphs and doc.paragraphs[0].style.name.startswith('Heading'):
                extracted_text += f"# {doc.paragraphs[0].text}\n\n"
                start_idx = 1
            else:
                start_idx = 0
            
            # Process paragraphs
            for para in doc.paragraphs[start_idx:]:
                # Handle headings
                if para.style.name.startswith('Heading 1'):
                    extracted_text += f"# {para.text}\n\n"
                elif para.style.name.startswith('Heading 2'):
                    extracted_text += f"## {para.text}\n\n"
                elif para.style.name.startswith('Heading 3'):
                    extracted_text += f"### {para.text}\n\n"
                # Handle lists
                elif para.style.name.startswith('List'):
                    extracted_text += f"- {para.text}\n"
                # Normal paragraphs
                elif para.text.strip():
                    # Check for bold, italic, etc.
                    formatted_text = ""
                    for run in para.runs:
                        text = run.text
                        if run.bold and run.italic:
                            text = f"***{text}***"
                        elif run.bold:
                            text = f"**{text}**"
                        elif run.italic:
                            text = f"*{text}*"
                        formatted_text += text
                    
                    extracted_text += f"{formatted_text}\n\n"
            
            # Process tables
            for table in doc.tables:
                # Create markdown table
                for i, row in enumerate(table.rows):
                    for cell in row.cells:
                        extracted_text += f"| {cell.text} "
                    extracted_text += "|\n"
                    
                    # Add separator after header row
                    if i == 0:
                        extracted_text += "|" + "---|" * len(row.cells) + "\n"
                
                extracted_text += "\n"
            
        elif textract_available:
            # Fallback to textract for other document types including .doc
            # Use textract to extract text directly from the file
            doc_text = textract.process(word_file).decode('utf-8')
            
            # Basic formatting for paragraphs
            paragraphs = doc_text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    extracted_text += f"{para.strip()}\n\n"
        else:
            return "Word document text extraction failed. Required libraries not installed."
        
        # Clean up the text
        extracted_text = clean_text(extracted_text)
        
        return extracted_text
            
    except Exception as e:
        logger.error(f"Error converting Word document to Markdown: {str(e)}")
        return f"Error extracting text from Word document: {str(e)}"

def convert_to_text(fileName: str, file_type: str) -> str:
    """
    Convert attachment content to plain text based on file type.
    
    Args:
        content: The content or text representation of the file
        file_type: The MIME type or file extension of the attachment
        
    Returns:
        Extracted text from the attachment
    """
    # If content is already provided as text, just return it
    if not fileName:
        return ""
    
    # Normalize file type to lowercase
    file_type = file_type.lower().strip()
    
    try:
        # Handle different file types based on MIME types or extensions
        # PDF files
        if "application/pdf" in file_type or file_type.endswith(".pdf"):
            return convert_pdf_to_markdown(fileName)
            
        # Word documents - use exact matching for MIME types and extension checking for simple types
        elif (file_type == "application/msword" or 
              file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or
              file_type.endswith(".docx") or file_type.endswith(".doc") or 
              file_type == "word"):
            return convert_word_to_markdown(fileName)
            
        # Excel files - use exact matching for MIME types and extension checking for simple types
        elif (file_type == "application/vnd.ms-excel" or
              file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or
              file_type.endswith(".xlsx") or file_type.endswith(".xls") or
              file_type == "excel"):
            return convert_excel_to_json(fileName)
            
        # Text files
        elif any(x in file_type for x in ["text/plain", "text/csv", "text/markdown", "text/tab-separated-values",
                                         "text/", "txt", "csv", "md", "tsv"]):
            try:
                with open(fileName, 'r', encoding='utf-8') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading text file {fileName}: {str(e)}")
                return f"Error reading file: {str(e)}"
             
        # PowerPoint files
        elif any(x in file_type for x in ["application/vnd.ms-powerpoint", 
                                         "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                         "powerpoint", "ppt", "pptx"]):
            return convert_ppt_to_markdown(fileName)
            
        # HTML content
        elif "text/html" in file_type or "html" in file_type:
            return html_to_markdown(fileName)
            
        # Default case - return content as is
        else:
            return ""
    
    except Exception as e:
        logger.error(f"Error extracting text from {file_type} attachment: {str(e)}")
        return ""

def base64_to_text(base64_string: str, encoding: str = 'utf-8') -> str:
    """
    Convert a base64-encoded string to plain text.
    
    Args:
        base64_string: The base64-encoded string to decode
        encoding: The text encoding to use (default: utf-8)
        
    Returns:
        Decoded plain text
    """
    if not base64_string:
        return ""
    
    try:
        # Remove any padding or whitespace that might affect decoding
        base64_string = base64_string.strip()
        
        # Handle potential URL-safe base64 encoding
        base64_string = base64_string.replace('-', '+').replace('_', '/')
        
        # Add padding if necessary
        padding = len(base64_string) % 4
        if padding:
            base64_string += '=' * (4 - padding)
            
        # Decode the base64 string to bytes
        decoded_bytes = base64.b64decode(base64_string)
        
        # Convert bytes to string using the specified encoding
        decoded_text = decoded_bytes.decode(encoding)
        
        return decoded_text
        
    except Exception as e:
        logger.error(f"Error decoding base64 string: {str(e)}")
        return f"Failed to decode base64 string: {str(e)}"
